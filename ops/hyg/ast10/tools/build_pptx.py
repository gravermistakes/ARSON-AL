#!/usr/bin/env python3
"""
Build a professional 12-slide PPTX for the OWASP Agentic Skills Top 10.

Slides:
  1.  Overview — all ten risks on one slide
  2-11. One slide per risk (AST01-AST10): issues vs mitigations diagram
  12. Summary — severity distribution + key takeaways

Content is parsed live from the astNN.md source files, so the deck stays in
sync with the documentation. The OWASP logo is placed on every slide.

Usage:
  python tools/build_pptx.py [--out dist/OWASP-Agentic-Skills-Top10.pptx]
"""
import argparse
import os
import re
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ---- palette (matches the site's visual Top 10 page) ----
CRIT  = RGBColor(0xC0, 0x25, 0x2B)
HIGH  = RGBColor(0xD9, 0x70, 0x1B)
MED   = RGBColor(0xB8, 0x8A, 0x00)
NAVY  = RGBColor(0x13, 0x39, 0x5C)
INK   = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x5B, 0x66, 0x75)
LINE  = RGBColor(0xD8, 0xDE, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREYBG = RGBColor(0xF7, 0xF9, 0xFC)
TEAL  = RGBColor(0x0F, 0x6E, 0x56)
ISSUE_BG = RGBColor(0xFC, 0xEC, 0xEC)
MIT_BG   = RGBColor(0xE3, 0xF3, 0xEC)

SEV_COLOR = {"Critical": CRIT, "High": HIGH, "Medium": MED, "Low": MED}


def sev_color(sev):
    return SEV_COLOR.get(sev, NAVY)


# ---------------------------------------------------------------- parsing
def smart_trunc(s, n=240):
    if len(s) <= n:
        return s
    cut = s[:n]
    dot = cut.rfind(". ")
    if dot >= 150:
        return cut[:dot + 1]
    sp = cut.rfind(" ")
    return cut[:sp if sp > 0 else n].rstrip(",;:— ") + "…"


def clean(text):
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\*\*([^*]*)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]*)\*", r"\1", text)
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    return re.sub(r"\s+", " ", text).strip()


def parse_ast(path):
    txt = open(path, encoding="utf-8").read()
    title = re.search(r"^title:\s*(.+)$", txt, re.M)
    title = title.group(1).strip().strip('"') if title else os.path.basename(path)
    m = re.match(r"(AST\d+)\s*[—\-:]\s*(.+)", title)
    ast_id, name = (m.group(1), m.group(2).strip()) if m else (title[:5], title)

    sev = re.search(r"\*\*Severity\*\*:\s*([A-Za-z]+)", txt)
    sev = sev.group(1) if sev else "Medium"

    desc = ""
    dm = re.search(r"##\s+Description\s*\n+(.+?)(?:\n\n|\n##)", txt, re.S)
    if dm:
        desc = clean(dm.group(1))

    issues = section_headings(txt, "Attack Scenarios")
    mits = mitigations(txt)
    return {
        "id": ast_id, "name": name, "sev": sev, "desc": desc,
        "issues": issues[:4], "mits": mits[:4],
    }


def section_block(txt, header):
    m = re.search(r"##\s+" + re.escape(header) + r"\s*\n(.*?)(?:\n##\s|\Z)", txt, re.S)
    return m.group(1) if m else ""


def section_headings(txt, header):
    block = section_block(txt, header)
    return [clean(h) for h in re.findall(r"^###\s+(.+)$", block, re.M)]


def mitigations(txt):
    block = section_block(txt, "Preventive Mitigations")
    out, seen = [], set()
    for line in re.findall(r"^\d+\.\s+(.+)$", block, re.M):
        b = re.search(r"\*\*(.+?)\*\*", line)
        label = clean(b.group(1)) if b else clean(line.split(":")[0])
        label = label.rstrip(".:")
        if label and label.lower() not in seen:
            seen.add(label.lower())
            out.append(label)
    return out


# ---------------------------------------------------------------- drawing helpers
def white_logo(src):
    """Return a path to a white-tinted copy of the (black) OWASP logo."""
    try:
        from PIL import Image
        img = Image.open(src).convert("RGBA")
        px = [(255, 255, 255, a) for (r, g, b, a) in img.getdata()]
        img.putdata(px)
        fd, p = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        img.save(p)
        return p
    except Exception:
        return src


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(slide, shape_type, x, y, w, h, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(s)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def set_text(shape, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             wrap=True, lpad=0.1):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("left", "right"):
        setattr(tf, "margin_" + m, Inches(lpad))
    for m in ("top", "bottom"):
        setattr(tf, "margin_" + m, Inches(0.03))
    if isinstance(runs, str):
        runs = [(runs, 14, INK, False)]
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, size, color, bold) in enumerate(runs):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tf


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb, runs, align=align, anchor=anchor, lpad=0.0)
    return tb


def add_logo(slide, on_dark, logo_black, logo_white, x, y, h):
    src = logo_white if on_dark else logo_black
    slide.shapes.add_picture(src, Inches(x), Inches(y), height=Inches(h))


VERSION = "v0.5"


def footer(slide, page):
    textbox(slide, 0.6, 7.06, 9, 0.3,
            [(f"OWASP Agentic Skills Top 10   ·   {VERSION}   ·   owasp.org/www-project-agentic-skills-top-10",
              9, MUTED, False)])
    textbox(slide, 11.8, 7.06, 1.0, 0.3, [(f"{page} / 12", 9, MUTED, False)],
            align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- slides
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def pill(slide, x, y, w, h, text, fg, bg):
    s = rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=bg)
    set_text(s, [(text, 11, fg, True)], align=PP_ALIGN.CENTER, lpad=0.05)
    return s


def overview_slide(prs, risks, logo_b, logo_w):
    s = blank(prs)
    rect(s, MSO_SHAPE.RECTANGLE, -0.05, -0.05, 13.44, 1.7, fill=NAVY)
    textbox(s, 0.6, 0.34, 10.5, 0.8, [("OWASP Agentic Skills Top 10", 34, WHITE, True)])
    textbox(s, 0.62, 1.06, 10.5, 0.45,
            [("The 10 most critical security risks for AI agent skills", 15,
              RGBColor(0xC7, 0xD3, 0xE6), False)])
    add_logo(s, True, logo_b, logo_w, 11.25, 0.5, 0.62)

    x_cols = [0.55, 6.85]
    w = 5.92
    y0, step, h = 1.98, 0.92, 0.78
    for i, r in enumerate(risks):
        col, row = divmod(i, 5)
        x = x_cols[col]
        y = y0 + row * step
        card = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=WHITE, line=LINE, line_w=1.0)
        rect(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, h, fill=sev_color(r["sev"]))
        textbox(s, x + 0.28, y + 0.12, 1.2, 0.5, [(r["id"], 15, sev_color(r["sev"]), True)],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 1.45, y + 0.12, w - 2.7, 0.55, [(r["name"], 14, INK, True)],
                anchor=MSO_ANCHOR.MIDDLE)
        pill(s, x + w - 1.18, y + 0.22, 1.02, 0.34, r["sev"].upper(), WHITE, sev_color(r["sev"]))

    ly = 6.66
    for i, (lab, c) in enumerate([("Critical", CRIT), ("High", HIGH), ("Medium", MED)]):
        lx = 0.6 + i * 1.7
        rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, 0.22, 0.22, fill=c)
        textbox(s, lx + 0.3, ly - 0.04, 1.3, 0.3, [(lab, 12, MUTED, False)])
    footer(s, 1)


def risk_slide(prs, r, page, logo_b, logo_w):
    s = blank(prs)
    c = sev_color(r["sev"])
    rect(s, MSO_SHAPE.RECTANGLE, -0.05, -0.05, 13.44, 1.32, fill=c)
    pill(s, 0.6, 0.18, 1.45, 0.34, r["sev"].upper(), c, WHITE)
    textbox(s, 0.58, 0.58, 10.0, 0.62, [(f"{r['id']} — {r['name']}", 27, WHITE, True)])
    add_logo(s, True, logo_b, logo_w, 11.3, 0.42, 0.5)

    if r["desc"]:
        textbox(s, 0.6, 1.45, 12.15, 0.72, [(smart_trunc(r["desc"]), 13, MUTED, False)])

    head_y = 2.3
    textbox(s, 0.6, head_y, 5.7, 0.34, [("ISSUES  /  ATTACK VECTORS", 14, CRIT, True)])
    textbox(s, 6.95, head_y, 5.7, 0.34, [("MITIGATIONS", 14, TEAL, True)])
    rect(s, MSO_SHAPE.RECTANGLE, 0.6, head_y + 0.4, 2.4, 0.03, fill=CRIT)
    rect(s, MSO_SHAPE.RECTANGLE, 6.95, head_y + 0.4, 2.4, 0.03, fill=TEAL)

    box_y0, bh, gap, bw = 2.92, 0.82, 0.12, 5.55
    for i, issue in enumerate(r["issues"]):
        y = box_y0 + i * (bh + gap)
        b = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, bw, bh, fill=ISSUE_BG, line=CRIT, line_w=1.0)
        set_text(b, [(issue, 13, RGBColor(0x79, 0x1F, 0x1F), True)], lpad=0.18)
    for i, mit in enumerate(r["mits"]):
        y = box_y0 + i * (bh + gap)
        b = rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.95, y, bw, bh, fill=MIT_BG, line=TEAL, line_w=1.0)
        set_text(b, [(mit, 13, RGBColor(0x08, 0x50, 0x41), True)], lpad=0.18)

    rows = max(len(r["issues"]), len(r["mits"]), 1)
    mid = box_y0 + (rows * (bh + gap) - gap) / 2 - 0.28
    arrow = rect(s, MSO_SHAPE.RIGHT_ARROW, 6.27, mid, 0.6, 0.56, fill=RGBColor(0xB4, 0xBE, 0xCC))
    no_shadow(arrow)
    footer(s, page)


def summary_slide(prs, risks, logo_b, logo_w):
    s = blank(prs)
    rect(s, MSO_SHAPE.RECTANGLE, -0.05, -0.05, 13.44, 1.32, fill=NAVY)
    textbox(s, 0.58, 0.34, 10.0, 0.62, [("Summary", 30, WHITE, True)])
    add_logo(s, True, logo_b, logo_w, 11.3, 0.42, 0.5)

    counts = {"Critical": 0, "High": 0, "Medium": 0}
    for r in risks:
        counts[r["sev"]] = counts.get(r["sev"], 0) + 1
    cards = [("Critical", counts.get("Critical", 0), CRIT),
             ("High", counts.get("High", 0), HIGH),
             ("Medium", counts.get("Medium", 0), MED)]
    cw, gap = 3.84, 0.3
    for i, (lab, n, c) in enumerate(cards):
        x = 0.7 + i * (cw + gap)
        rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.65, cw, 1.45, fill=GREYBG, line=LINE, line_w=1.0)
        rect(s, MSO_SHAPE.RECTANGLE, x, 1.65, 0.14, 1.45, fill=c)
        textbox(s, x + 0.35, 1.74, cw - 0.5, 0.9, [(str(n), 40, c, True)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 1.45, 1.74, cw - 1.6, 1.3,
                [(f"{lab}-severity\nrisks", 16, INK, True)], anchor=MSO_ANCHOR.MIDDLE)

    textbox(s, 0.7, 3.45, 12, 0.4, [("Key takeaways", 18, NAVY, True)])
    takeaways = [
        "Skills run with the host agent's full privileges — one malicious or injected skill can reach credentials, files, and shell.",
        "Prompt injection amplifies every category: untrusted input and tool output can drive autonomous, over-privileged actions.",
        "Defenses compose — signing + provenance + least-privilege manifests + sandboxing + scanning + governance, not any one control.",
        "There is no universal skill format, so security metadata is lost when skills are ported across platforms.",
    ]
    y = 3.95
    for t in takeaways:
        rect(s, MSO_SHAPE.OVAL, 0.75, y + 0.08, 0.16, 0.16, fill=NAVY)
        textbox(s, 1.1, y, 11.4, 0.62, [(t, 14, INK, False)])
        y += 0.72

    rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 6.78, 12.0, 0.0, fill=None, line=LINE)
    textbox(s, 0.7, 6.62, 12, 0.35,
            [("Full details, checklist, and per-risk pages: ", 12, MUTED, False),
             ("owasp.org/www-project-agentic-skills-top-10", 12, NAVY, True)])
    footer(s, 12)


def build(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    logo_b = os.path.join(REPO, "assets", "images", "owasp-logo.png")
    logo_w = white_logo(logo_b) if os.path.exists(logo_b) else logo_b
    if not os.path.exists(logo_b):
        print("WARNING: OWASP logo not found at", logo_b, file=sys.stderr)
        logo_b = logo_w = None

    risks = []
    for i in range(1, 11):
        p = os.path.join(REPO, f"ast{i:02d}.md")
        if os.path.exists(p):
            risks.append(parse_ast(p))
    if len(risks) != 10:
        print(f"WARNING: parsed {len(risks)} risks (expected 10)", file=sys.stderr)

    overview_slide(prs, risks, logo_b, logo_w)
    for i, r in enumerate(risks):
        risk_slide(prs, r, i + 2, logo_b, logo_w)
    summary_slide(prs, risks, logo_b, logo_w)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path} ({len(prs.slides._sldIdLst)} slides)")


def main():
    global VERSION
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(REPO, "dist", "OWASP-Agentic-Skills-Top10.pptx"))
    ap.add_argument("--doc-version", dest="ver", default=VERSION)
    args = ap.parse_args()
    VERSION = args.ver
    build(args.out)


if __name__ == "__main__":
    main()
